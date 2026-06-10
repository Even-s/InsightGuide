#!/usr/bin/env python3
"""
Test script for Milestone 2: AI Slide Analysis + Topic Card Generation

This script tests the complete pipeline:
1. Upload PPTX file
2. Monitor processing status
3. Retrieve generated topic cards
4. Verify results
"""

import requests
import time
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# Configuration
API_BASE_URL = "http://localhost:8001/api"
TEST_PPTX_PATH = Path("/tmp/test_presentation.pptx")


def create_test_pptx():
    """Create a simple test PPTX file with meaningful content."""
    print("📝 Creating test PPTX file...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Introduction to Machine Learning
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide1.shapes.title.text = "Introduction to Machine Learning"

    content = slide1.placeholders[1].text_frame
    content.text = "Machine Learning Overview"

    p = content.add_paragraph()
    p.text = "• Supervised Learning: Learning from labeled data"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Unsupervised Learning: Finding patterns in unlabeled data"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Reinforcement Learning: Learning through trial and error"
    p.level = 1

    p = content.add_paragraph()
    p.text = "Key algorithms include neural networks, decision trees, and SVMs"
    p.level = 1

    # Slide 2: Neural Networks Deep Dive
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Neural Networks Architecture"

    content = slide2.placeholders[1].text_frame
    content.text = "Core Components"

    p = content.add_paragraph()
    p.text = "• Input Layer: Receives raw data features"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Hidden Layers: Process information through activation functions"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Output Layer: Produces final predictions"
    p.level = 1

    p = content.add_paragraph()
    p.text = "Training uses backpropagation and gradient descent"
    p.level = 1

    # Slide 3: Practical Applications
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Real-World Applications"

    content = slide3.placeholders[1].text_frame
    content.text = "ML in Industry"

    p = content.add_paragraph()
    p.text = "• Computer Vision: Image recognition, object detection"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Natural Language Processing: Chatbots, translation, sentiment analysis"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Recommendation Systems: Netflix, Amazon, Spotify"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Autonomous Vehicles: Self-driving cars use multiple ML techniques"
    p.level = 1

    prs.save(TEST_PPTX_PATH)
    print(f"✅ Test PPTX created at {TEST_PPTX_PATH}")
    return TEST_PPTX_PATH


def upload_deck(file_path):
    """Upload PPTX file to the API."""
    print(f"\n📤 Uploading deck: {file_path}")

    with open(file_path, 'rb') as f:
        files = {'file': (file_path.name, f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
        data = {'title': 'Test Presentation - Machine Learning'}

        response = requests.post(
            f"{API_BASE_URL}/decks/",
            files=files,
            data=data
        )

    if response.status_code != 201:
        print(f"❌ Upload failed: {response.status_code}")
        print(response.text)
        return None

    deck_data = response.json()
    print(f"✅ Deck uploaded successfully")
    print(f"   Deck ID: {deck_data['id']}")
    print(f"   Status: {deck_data['status']}")
    return deck_data['id']


def monitor_status(deck_id, max_wait=180, check_interval=3):
    """Monitor deck processing status until completion."""
    print(f"\n⏳ Monitoring deck status (max {max_wait}s)...")

    start_time = time.time()
    previous_status = None

    while time.time() - start_time < max_wait:
        response = requests.get(f"{API_BASE_URL}/decks/{deck_id}/status")

        if response.status_code != 200:
            print(f"❌ Status check failed: {response.status_code}")
            return None

        status_data = response.json()
        current_status = status_data['status']

        if current_status != previous_status:
            print(f"   📊 Status: {current_status} - {status_data['message']}")
            previous_status = current_status

        if current_status == 'analyzed':
            elapsed = time.time() - start_time
            print(f"✅ Analysis complete in {elapsed:.1f}s")
            return current_status

        if current_status == 'failed':
            print(f"❌ Processing failed")
            return current_status

        time.sleep(check_interval)

    print(f"⚠️  Timeout after {max_wait}s. Last status: {previous_status}")
    return previous_status


def get_analysis_results(deck_id):
    """Retrieve and display analysis results."""
    print(f"\n📊 Retrieving analysis results...")

    response = requests.get(f"{API_BASE_URL}/decks/{deck_id}/analysis")

    if response.status_code != 200:
        print(f"❌ Failed to get analysis: {response.status_code}")
        print(response.text)
        return None

    analysis = response.json()

    print(f"✅ Analysis Results:")
    print(f"   Deck ID: {analysis['deck_id']}")
    print(f"   Status: {analysis['status']}")
    print(f"   Total Slides: {len(analysis['slides'])}")
    print(f"   Total Topic Cards: {analysis['topic_cards_count']}")

    print(f"\n📑 Slides:")
    for slide in analysis['slides']:
        print(f"   • Slide {slide['page_number']}: {slide['title']}")
        print(f"     - Image: {slide['image_url']}")
        print(f"     - Topic Cards: {slide['topic_cards_count']}")
        if slide.get('ai_summary'):
            print(f"     - Summary: {slide['ai_summary'][:100]}...")

    return analysis


def get_topic_cards(deck_id):
    """Retrieve and display topic cards."""
    print(f"\n🎯 Retrieving topic cards...")

    response = requests.get(f"{API_BASE_URL}/topic-cards/deck/{deck_id}")

    if response.status_code != 200:
        print(f"❌ Failed to get topic cards: {response.status_code}")
        print(response.text)
        return None

    cards = response.json()

    print(f"✅ Found {len(cards)} topic cards")

    for i, card in enumerate(cards, 1):
        print(f"\n   Card {i}: {card['title']}")
        print(f"   - Topic Type: {card['topicType']}")
        print(f"   - Importance: {card['importance']}")
        print(f"   - Slide: {card['slideId']}")

        # Coverage rule is nested
        coverage_rule = card.get('coverageRule', {})

        if coverage_rule.get('semanticAnchors'):
            anchors = coverage_rule['semanticAnchors']
            print(f"   - Semantic Anchors ({len(anchors)}): {anchors[0][:80]}...")

        if coverage_rule.get('expectedKeywords'):
            keywords = coverage_rule['expectedKeywords']
            print(f"   - Expected Keywords ({len(keywords)}): {', '.join(keywords[:5])}...")

        if coverage_rule.get('mustMentionFacts'):
            facts = coverage_rule['mustMentionFacts']
            print(f"   - Must Mention Facts: {len(facts)} facts")

        if card.get('suggestedScript'):
            script = card['suggestedScript']
            print(f"   - Suggested Script: {script[:100]}...")

    return cards


def verify_results(analysis, topic_cards):
    """Verify the results meet expectations."""
    print(f"\n✅ Verification Results:")

    checks = []

    # Check 1: Slides were created
    if len(analysis['slides']) >= 3:
        print(f"   ✓ Slides created: {len(analysis['slides'])} slides")
        checks.append(True)
    else:
        print(f"   ✗ Expected at least 3 slides, got {len(analysis['slides'])}")
        checks.append(False)

    # Check 2: Topic cards were generated
    if len(topic_cards) > 0:
        print(f"   ✓ Topic cards generated: {len(topic_cards)} cards")
        checks.append(True)
    else:
        print(f"   ✗ No topic cards generated")
        checks.append(False)

    # Check 3: Topic cards have required fields
    if topic_cards:
        card = topic_cards[0]
        required_fields = ['title', 'topicType', 'importance', 'coverageRule']
        missing = [f for f in required_fields if not card.get(f)]

        if not missing:
            print(f"   ✓ Topic cards have required fields")
            checks.append(True)
        else:
            print(f"   ✗ Topic cards missing fields: {missing}")
            checks.append(False)

    # Check 4: Coverage rules are properly structured
    if topic_cards:
        card = topic_cards[0]
        coverage_rule = card.get('coverageRule', {})
        has_anchors = coverage_rule.get('semanticAnchors') and len(coverage_rule['semanticAnchors']) > 0
        has_keywords = 'expectedKeywords' in coverage_rule
        has_facts = 'mustMentionFacts' in coverage_rule
        has_thresholds = coverage_rule.get('thresholds') and 'covered' in coverage_rule['thresholds']

        if has_anchors and has_keywords and has_facts and has_thresholds:
            print(f"   ✓ Coverage rules properly defined")
            checks.append(True)
        else:
            print(f"   ✗ Coverage rules incomplete (anchors:{has_anchors}, keywords:{has_keywords}, facts:{has_facts}, thresholds:{has_thresholds})")
            checks.append(False)
    else:
        checks.append(False)

    # Check 5: Suggested scripts exist
    if topic_cards:
        cards_with_scripts = [c for c in topic_cards if c.get('suggestedScript')]
        if cards_with_scripts:
            print(f"   ✓ Suggested scripts generated ({len(cards_with_scripts)}/{len(topic_cards)} cards)")
            checks.append(True)
        else:
            print(f"   ✗ No suggested scripts found")
            checks.append(False)
    else:
        checks.append(False)

    # Summary
    passed = sum(checks)
    total = len(checks)
    print(f"\n{'='*60}")
    print(f"   Checks Passed: {passed}/{total}")

    if passed == total:
        print(f"   🎉 ALL TESTS PASSED!")
    else:
        print(f"   ⚠️  Some tests failed - review above for details")
    print(f"{'='*60}")

    return passed == total


def main():
    """Run the complete Milestone 2 test."""
    print("=" * 60)
    print("   MILESTONE 2 TEST: AI Slide Analysis + Topic Card Generation")
    print("=" * 60)

    try:
        # Step 1: Create test PPTX
        test_file = create_test_pptx()

        # Step 2: Upload deck
        deck_id = upload_deck(test_file)
        if not deck_id:
            print("\n❌ Test failed: Could not upload deck")
            return False

        # Step 3: Monitor status
        final_status = monitor_status(deck_id, max_wait=180)
        if final_status != 'analyzed':
            print(f"\n❌ Test failed: Deck did not reach 'analyzed' status (got: {final_status})")
            return False

        # Step 4: Get analysis results
        analysis = get_analysis_results(deck_id)
        if not analysis:
            print("\n❌ Test failed: Could not retrieve analysis")
            return False

        # Step 5: Get topic cards
        topic_cards = get_topic_cards(deck_id)
        if not topic_cards:
            print("\n⚠️  Warning: No topic cards found, but continuing verification")
            topic_cards = []

        # Step 6: Verify results
        success = verify_results(analysis, topic_cards)

        print(f"\n{'='*60}")
        if success:
            print("   ✅ MILESTONE 2 TEST: PASSED")
        else:
            print("   ❌ MILESTONE 2 TEST: FAILED")
        print(f"{'='*60}")

        return success

    except Exception as e:
        print(f"\n❌ Test failed with exception: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)
